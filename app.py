import streamlit as st

st.set_page_config(layout='wide', initial_sidebar_state='expanded', page_title="Strategic Synthesis AI")

import os
import streamlit as st
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from langchain_ollama import ChatOllama
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, AIMessage
import pdfplumber
from langchain_community.chat_message_histories import StreamlitChatMessageHistory
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.runnables.history import RunnableWithMessageHistory
from src.data.rag import get_rich_case_studies
from langchain_google_genai import ChatGoogleGenerativeAI

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
import tempfile
import random

# Initialize session states
if 'welcome_complete' not in st.session_state:
    st.session_state.welcome_complete = False
if 'case_studies_fetched' not in st.session_state:
    st.session_state.case_studies_fetched = False
    st.session_state.context_prompt = ""

ENABLE_DECORATIVE_CIRCLES = False
ENABLE_TITLE_BAR_PATTERNS = False
ENABLE_BACKGROUND_PATTERNS = True  # You can toggle this too

def get_rgb_components(color: RGBColor):
    color_str = str(color)  # e.g., 'FF6699'
    r = int(color_str[0:2], 16)
    g = int(color_str[2:4], 16)
    b = int(color_str[4:6], 16)
    return r, g, b

# --- Enhanced Style Constants ---
TITLE_FONT_SIZE = Pt(44)
SUBTITLE_FONT_SIZE = Pt(28)
HEADING_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(20)

# Color schemes - you can switch between these
COLOR_SCHEMES = {
    "blue": {
        "primary": RGBColor(0, 102, 204),      # Main accent color
        "secondary": RGBColor(41, 128, 185),   # Secondary accent
        "accent": RGBColor(52, 152, 219),      # Highlight color
        "text": RGBColor(44, 62, 80),          # Text color
        "light_text": RGBColor(255, 255, 255), # Light text color
        "background": RGBColor(240, 240, 240), # Slide background
        "alt_background": RGBColor(245, 250, 255), # Alternative background
        "chart_colors": [
            RGBColor(41, 128, 185),
            RGBColor(26, 188, 156),
            RGBColor(46, 204, 113),
            RGBColor(241, 196, 15)
        ]
    },
    "green": {
        "primary": RGBColor(46, 139, 87),     
        "secondary": RGBColor(60, 179, 113),  
        "accent": RGBColor(76, 187, 23),      
        "text": RGBColor(44, 62, 80),         
        "light_text": RGBColor(255, 255, 255),
        "background": RGBColor(240, 245, 240),
        "alt_background": RGBColor(245, 255, 245),
        "chart_colors": [
            RGBColor(46, 139, 87),
            RGBColor(115, 169, 80),
            RGBColor(39, 174, 96),
            RGBColor(241, 196, 15)
        ]
    },
    "purple": {
        "primary": RGBColor(142, 68, 173),     
        "secondary": RGBColor(155, 89, 182),   
        "accent": RGBColor(187, 143, 206),      
        "text": RGBColor(44, 62, 80),          
        "light_text": RGBColor(255, 255, 255), 
        "background": RGBColor(245, 240, 250), 
        "alt_background": RGBColor(250, 245, 255),
        "chart_colors": [
            RGBColor(142, 68, 173),
            RGBColor(155, 89, 182),
            RGBColor(26, 188, 156),
            RGBColor(241, 196, 15)
        ]
    }
}

# Choose a color scheme (can be randomized or user-selected)
COLORS = COLOR_SCHEMES["blue"]  # Default color scheme

def clean_financial_slide(prs):
    """
    Removes non-text decorative shapes from the slide titled 'Financial Projections'.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text is not None:
                if "financial projections" in shape.text.strip().lower():
                    shapes_to_remove = [
                        s for s in slide.shapes
                        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not s.has_text_frame
                    ]
                    for s in shapes_to_remove:
                        s._element.getparent().remove(s._element)
                    return
            
# --- File Text Extractor Functions ---
def extract_text(file):
    if file.type == "application/pdf":
        return extract_pdf_text(file)
    elif file.type == "text/plain":
        return file.read().decode("utf-8")
    elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return extract_docx_text(file)
    else:
        return f"❌ Unsupported file type: {file.type}"

def extract_pdf_text(file):
    text = ""
    with fitz.open(stream=file.read(), filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_docx_text(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

# --- Background Pattern Generators ---
def add_dotted_background(slide, color=None, density=120):
    """Add dotted pattern background to a slide"""
    if color is None:
        color = COLORS["accent"]
    
    # Create small dots scattered across the slide
    for _ in range(density):
        x = random.uniform(0, 10)
        y = random.uniform(0, 7.5)
        dot_size = random.uniform(0.02, 0.05)  # Small sized dots
        
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, 
            Inches(x), 
            Inches(y), 
            Inches(dot_size), 
            Inches(dot_size)
        )
        dot.fill.solid()
        
        # Create a slightly transparent version of the color
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        # Set transparency (0-100000, where 100000 is fully transparent)
        dot.fill.transparency = random.randint(60000, 85000)

def add_wave_background(slide, color1=None, color2=None):
    """Add a subtle wave pattern to slide background"""
    if color1 is None:
        color1 = COLORS["primary"]
    if color2 is None:
        color2 = COLORS["secondary"]
    
    # Create a few wave shapes
    for i in range(3):
        height = random.uniform(0.8, 1.2)
        y_pos = random.uniform(-0.5, 7)
        
        wave = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.WAVE, 
            Inches(-0.5), 
            Inches(y_pos), 
            Inches(11), 
            Inches(height)
        )
        wave.fill.solid()
        wave.fill.fore_color.rgb = color1 if i % 2 == 0 else color2
        wave.line.fill.background()
        wave.fill.transparency = random.randint(70000, 90000)

def add_geometric_background(slide):
    """Add geometric shapes as background elements"""
    shapes = [
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        MSO_AUTO_SHAPE_TYPE.OVAL,
        MSO_AUTO_SHAPE_TYPE.DIAMOND,
        MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM
        ]
    
    for _ in range(15):
        shape_type = random.choice(shapes)
        size = random.uniform(0.5, 2.0)
        x = random.uniform(-0.5, 10)
        y = random.uniform(-0.5, 7.5)
        
        shape = slide.shapes.add_shape(
            shape_type, 
            Inches(x), 
            Inches(y), 
            Inches(size), 
            Inches(size)
        )
        
        # Pick a color from the color scheme
        color_keys = ["primary", "secondary", "accent"]
        color = COLORS[random.choice(color_keys)]
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.fill.transparency = random.randint(80000, 95000)  # Very transparent

def add_gradient_background(slide, direction="vertical"):
    """
    This simulates a gradient effect by adding multiple semi-transparent rectangles
    """
    steps = 15  # Number of gradient steps
    width = Inches(10)
    height = Inches(7.5)
    
    # Choose two colors for the gradient
    start_color = COLORS["primary"]
    end_color = COLORS["secondary"]

    r1, g1, b1 = get_rgb_components(start_color)
    r2, g2, b2 = get_rgb_components(end_color)

    for i in range(steps):
        # Calculate interpolated color
        ratio = i / steps
        r = int(r1 * (1 - ratio) + r2 * ratio)
        g = int(g1 * (1 - ratio) + g2 * ratio)
        b = int(b1 * (1 - ratio) + b2 * ratio)
        
        if direction == "vertical":
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0),
                Inches(7.5 * i / steps),
                width,
                Inches(7.5 / steps + 0.1)  # Slight overlap to avoid gaps
            )
        else:  # horizontal
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(10 * i / steps),
                Inches(0),
                Inches(10 / steps + 0.1),  # Slight overlap
                height
            )
            
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()
        rect.fill.transparency = 50000  # Semi-transparent

# --- Title Slide ---
def add_styled_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    add_geometric_background(slide)

    center_overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    )
    center_overlay.fill.solid()
    center_overlay.fill.fore_color.rgb = COLORS["background"]
    center_overlay.line.fill.background()
    center_overlay.fill.transparency = 20000

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.8), Inches(2), Inches(0.15), Inches(3.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS["primary"]
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(2.1), Inches(7.4), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(Inches(1.6), Inches(3.5), Inches(7.4), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    p = subtitle_tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["secondary"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    if ENABLE_DECORATIVE_CIRCLES:
        for i in range(3):
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(9 - i*0.3), Inches(5.5 + i*0.2), Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS["primary"] if i % 2 == 0 else COLORS["accent"]
            circle.line.fill.background()

# --- Section Divider Slide ---
def add_section_divider_slide(prs, title):
    """Creates an impactful section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add gradient background
    add_gradient_background(slide, direction="horizontal")
    
    # Add a semi-transparent overlay for better text visibility
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS["primary"]
    overlay.line.fill.background()
    overlay.fill.transparency = 40000  # Somewhat transparent
    
    # Large, bold title text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = title.upper()  # Uppercase for impact
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Add decorative line
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(4),
        Inches(4.5),
        Inches(2),
        Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["light_text"]
    line.line.fill.background()

# --- Enhanced Bullet Slide with Optional Infographic ---
def add_bullet_slide(prs, title, bullets, infographic=None, background_style="dots"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set base background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    
    # Add background pattern based on style
    if ENABLE_BACKGROUND_PATTERNS:
        if background_style == "dots":
            add_dotted_background(slide)
        elif background_style == "waves":
            add_wave_background(slide)
        elif background_style == "geometric":
            add_geometric_background(slide)
    
    # Enhanced title bar
    title_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 
        Inches(0), 
        Inches(0), 
        Inches(10), 
        Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS["primary"]
    title_shape.line.fill.background()
    
    # Add subtle pattern to title bar
    if ENABLE_TITLE_BAR_PATTERNS:
        for i in range(5):
            x = random.uniform(0, 9)
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(x),
                Inches(random.uniform(0.1, 0.9)),
                Inches(0.3),
                Inches(0.3)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS["secondary"]
            circle.line.fill.background()
            circle.fill.transparency = 60000
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    p.font.name = "Calibri"
    
    # Content area with slight shadow effect
    content_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(1.3),
        Inches(6.8),
        Inches(5.5)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_bg.line.color.rgb = RGBColor(220, 220, 220)
    
    # Content with custom bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        para = tf.add_paragraph() if i != 0 else tf.paragraphs[0]
        
        # Alternate bullet styles for visual interest
        if i % 2 == 0:
            para.text = f"► {bullet}"  # Triangle bullet
        else:
            para.text = f"• {bullet}"  # Round bullet
            
        para.font.size = BODY_FONT_SIZE
        para.font.color.rgb = COLORS["text"]
        para.font.name = "Calibri"
        
        # Add spacing between bullets
        if i > 0:
            para.space_before = Pt(12)

    # Add infographic if provided
    if infographic:
        infographic(slide)

# --- Enhanced SWOT Infographic Slide ---
def add_swot_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    add_dotted_background(slide)

    tb = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = COLORS["primary"]
    tb.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "SWOT Analysis"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["light_text"]
    title_p.font.name = "Calibri"

    labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    coords = [(0.6, 1.4), (5.1, 1.4), (0.6, 4.0), (5.1, 4.0)]
    colors = [
        RGBColor(76, 175, 80),
        RGBColor(244, 67, 54),
        RGBColor(33, 150, 243),
        RGBColor(255, 152, 0)
    ]

    for i, label in enumerate(labels):
        x, y = coords[i]

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(2.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = colors[i]
        box.line.width = Pt(2.5)

        header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors[i]
        header.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.2), Inches(0.5))
        label_tf = label_box.text_frame
        p = label_tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        text_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.7), Inches(3.8), Inches(1.4)
        )
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        placeholder_p = text_tf.paragraphs[0]
        placeholder_p.text = f"Add {label.lower()} here"
        placeholder_p.font.italic = True
        placeholder_p.font.size = Pt(14)
        placeholder_p.font.color.rgb = RGBColor(120, 120, 120)

# --- Roadmap Slide ---
def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    add_dotted_background(slide)

    tb = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = COLORS["primary"]
    tb.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "Implementation Roadmap"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["light_text"]
    title_p.font.name = "Calibri"

    timeline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(3), Inches(8), Inches(0.1)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = COLORS["text"]
    timeline.line.fill.background()

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(9), Inches(2.9), Inches(0.5), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["text"]
    arrow.line.fill.background()

    phases = [
        {"name": "Planning", "duration": "0-3 months", "color": COLORS["chart_colors"][0]},
        {"name": "Rollout", "duration": "4-6 months", "color": COLORS["chart_colors"][1]},
        {"name": "Expansion", "duration": "7-12 months", "color": COLORS["chart_colors"][2]}
    ]

    for i, phase in enumerate(phases):
        x = 2 + i * 3

        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x), Inches(3), Inches(0.8), Inches(0.8)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = phase["color"]
        node.line.fill.background()
        p = node.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["light_text"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        lbl = slide.shapes.add_textbox(
            Inches(x - 0.6), Inches(4), Inches(2), Inches(0.6)
        )
        p2 = lbl.text_frame.paragraphs[0]
        p2.text = phase["name"]
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLORS["text"]
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER

        dur = slide.shapes.add_textbox(
            Inches(x - 0.6), Inches(4.4), Inches(2), Inches(0.5)
        )
        p3 = dur.text_frame.paragraphs[0]
        p3.text = phase["duration"]
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = COLORS["secondary"]
        p3.font.name = "Calibri"
        p3.alignment = PP_ALIGN.CENTER

        act_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x - 0.75), Inches(1.8), Inches(2.3), Inches(0.9)
        )
        act_box.fill.solid()
        act_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        act_box.line.color.rgb = phase["color"]

        act_text = slide.shapes.add_textbox(
            Inches(x - 0.65), Inches(1.9), Inches(2.1), Inches(0.7)
        )
        act_tf = act_text.text_frame
        act_tf.word_wrap = True
        p4 = act_tf.paragraphs[0]
        p4.text = "Key Activities:"
        p4.font.size = Pt(12)
        p4.font.bold = True
        p4.font.color.rgb = COLORS["text"]
        p5 = act_tf.add_paragraph()
        p5.text = "• Activity 1\n• Activity 2"
        p5.font.size = Pt(10)
        p5.font.color.rgb = COLORS["text"]

# --- Financial Chart Slide ---
def add_financial_chart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "Financial Projections"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["primary"]
    title_p.font.name = "Calibri"

    # Chart bars (simple visual simulation)
    items = [
        ("Initial Investment", 2.5, COLORS["chart_colors"][0]),
        ("Expected Revenue", 4.2, COLORS["chart_colors"][1]),
        ("Operational Costs", 2.0, COLORS["chart_colors"][2]),
        ("Break-even", 3.1, COLORS["chart_colors"][3])
    ]

    for i, (label, height, color) in enumerate(items):
        x = 1.0 + i * 2.1
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(5.5 - height),
            Inches(1),
            Inches(height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(5.7), Inches(1.4), Inches(0.5))
        label_tf = label_box.text_frame
        p = label_tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["text"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# --- Generate Full Styled PPT ---
def generate_styled_pptx():
    prs = Presentation()

    add_styled_title_slide(prs, "Business Growth Strategy Proposal", "Prepared by Strategic Synthesis AI")

    add_section_divider_slide(prs, "Introduction")
    add_bullet_slide(prs, "Company Overview", [
        "Founded in 2020 with a vision to democratize access to AI.",
        "Headquartered in New York with global outreach.",
        "Team of 50+ AI experts and software engineers."
    ])

    add_section_divider_slide(prs, "Business Landscape")
    add_bullet_slide(prs, "Market Understanding", [
        "Target market includes mid-sized B2B SaaS providers.",
        "Market size is projected to reach $5B by 2027.",
        "Customer pain points revolve around data silos and integration."
    ])
    add_swot_slide(prs)

    add_section_divider_slide(prs, "Strategy & Roadmap")
    add_bullet_slide(prs, "Strategic Objectives", [
        "Expand to 3 new international markets.",
        "Achieve $10M ARR by Q4 2026.",
        "Establish partnerships with key ecosystem players."
    ])
    add_roadmap_slide(prs)

    add_section_divider_slide(prs, "Financial Outlook")
    add_financial_chart_slide(prs)
    clean_financial_slide(prs)

    add_section_divider_slide(prs, "Outcomes & Next Steps")
    add_bullet_slide(prs, "Expected Impact", [
        "Revenue growth by 300% over 2 years.",
        "Improved brand positioning in AI infrastructure sector.",
        "Sustainable operational margin with lean team."
    ])

    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(output_file.name)
    return output_file.name

# Model Settings
model_provider = os.getenv("MODEL_PROVIDER", "ollama")  # 'ollama' or 'openai'
llm_model = os.getenv("LLM_MODEL", "Gemma3")  # Default for ollama

def get_llm():
    """Initialize LLM model based on environment settings."""
    llm = None

    if model_provider.lower() == "openai":
        llm = ChatOpenAI(
            model_name=os.getenv("LLM_MODEL", "gpt-4"),
            temperature=0.1
        )
    elif model_provider.lower() == "google":
        llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash-001",
        temperature=0,
        max_tokens=None,
        timeout=5,
        max_retries=2,
        # other params...
    )
    else:  # Default to ollama
        llm = ChatOllama(
            model=llm_model,
            temperature=0.1
        )

    return llm

# File Text Extractor Functions
def extract_text(file):
    if file.type == "application/pdf":
        return extract_pdf_text(file)
    elif file.type == "text/plain":
        return file.read().decode("utf-8")
    else:
        return f"❌ Unsupported file type: {file.type}"

def extract_pdf_text(file):
    text = ""
    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text

# Function to handle starting the app
def start_app():
    st.session_state.welcome_complete = True

# Welcome Screen
if not st.session_state.welcome_complete:
    col1, col2, col3 = st.columns([1, 3, 1])
    
    with col2:
        st.markdown(
            """
            <div style="background-color: #222; padding: 30px; border-radius: 10px; margin-bottom: 20px; text-align: center;">
                <div style="display: flex; align-items: center; justify-content: center; margin-bottom: 10px;">
                    <span style="color: #FFB627; font-size: 40px; margin-right: 10px;">&#9733;</span>
                    <h1 style="color: white; margin: 0;">Strategic Synthesis AI</h1>
                </div>
            </div>
            """, 
            unsafe_allow_html=True
        )
        
        st.markdown(
            """
            <h1 style="text-align: center; font-size: 2.5rem; margin-bottom: 30px;">
            Your friendly neighborhood business advisor
            </h1>
            """, 
            unsafe_allow_html=True
        )
        
        st.markdown(
            """
            <p style="text-align: center; font-size: 1.2rem; margin-bottom: 40px; color: #555;">
            Strategic Synthesis AI analyzes your business data, generates comprehensive strategy plans, 
            and builds tailored presentations for different stakeholders using the power of AI
            </p>
            """, 
            unsafe_allow_html=True
        )
        
        col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
        with col_btn2:
            if st.button("Lets get workin'!", use_container_width=True, type="primary"):
                start_app()
                st.rerun()

# Main Application
else:
    st.title("Strategic Synthesis AI")

    # File uploader
    uploaded_files = st.file_uploader("Attach anything that will help me understand your business 😄", accept_multiple_files=True)
    attached_text = ""  # Store all uploaded content here

    if uploaded_files:
        for file in uploaded_files:
            extracted = extract_text(file)
            attached_text += f"\n\n--- File: {file.name} ---\n{extracted}"
        
        # Only fetch case studies if we have new content and haven't fetched them yet
        if attached_text and not st.session_state.case_studies_fetched:
            with st.spinner('Tapping into the infinite wisdom of universe...'):
                st.session_state.context_prompt = str(get_rich_case_studies(attached_text)).replace("{", "{{").replace("}", "}}")
                st.session_state.case_studies_fetched = True
                st.success("Case studies retrieved successfully!")

    llm = get_llm()

    # Use the stored context prompt from session state
    context_prompt = st.session_state.context_prompt

    # Prompt Template (inject file content here)
    system_prompt = f"""
    You are a world-class Strategic Business Advisor helping businesses across industries design effective business strategies to support growth, innovation, and long-term success.

    Carefully analyze the content provided by the user (if it is provided) to understand their current business model, challenges, or goals.

    You should explicitly refer to the reference case studies provided by the system to craft thoughtful, tailored, and actionable responses to the user's questions.
    It is very important to provide reference http URLs from the Relevant Case Studies in your response

    Analyze the following key factors for your reponse:-
    1. 	Business Context Industry specifics (e.g., tech, FMCG, healthcare) -Market structure (monopoly, oligopoly, fragmented, etc.), Regulatory environment (local/global, highly regulated or not), Stage of business lifecycle (startup, growth, maturity, decline)
    2. 	Strategic Intent Vision/Mission alignment - Growth objectives (scale, profit, market leadership, innovation), Geographic goals (domestic focus vs global expansion), Long-term vs. short-term priorities
    3. 	Key Capability Inputs Core competencies (e.g., R&D, brand strength, distribution) - Technology maturity, Talent & leadership, Operational infrastructure
    4. 	Customer Dimensions Target segment behavior - Customer jobs to be done, Channel preferences (D2C, retail, B2B), Value perception and willingness to pay
    5. 	Competitive Forces Rivalry intensity - Barriers to entry, Threat of substitutes, Supplier/buyer power (Porter's Five Forces), Innovation speed in the industry
    6. 	Strategic Options Spectrum Growth levers (market penetration, product dev, M&A, diversification) - Business models (B2B/B2C, SaaS, subscription, platform), Differentiation methods (price, innovation, service, brand), Focus areas (niche vs mass market)
    7. 	Risk & Resilience Metrics Financial risk tolerance, Operational risk (supply chain fragility), Market volatility exposure, Crisis adaptability (e.g., COVID learnings)
    8. 	Measurement and Governance Key Performance Indicators (KPIs) -Feedback loops, Decision accountability, Scalability of strategy

    Based upon this analysis, what are the gaps in the existing business and develop detailed 3 month, 6 month and 1 year plan y first asking user what they want to know.

    Do not make up information or assumptions. If the user content is insufficient to fully answer a question, clearly say so and suggest what additional information would help.

    Be clear, insightful, and professional — your goal is to act as a trusted strategic advisor.

    --- Attached User Content ---
    {attached_text}

    --- Relevant Case Studies ---
    {context_prompt}

    """

    # Set up memory
    msgs = StreamlitChatMessageHistory(key="langchain_messages")

    prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system_prompt),
            MessagesPlaceholder(variable_name="history"),
            ("human", "{question}"),
        ]
    )

    chain = prompt | llm
    chain_with_history = RunnableWithMessageHistory(
        chain,
        lambda session_id: msgs,
        input_messages_key="question",
        history_messages_key="history",
    )

    # Render current messages from StreamlitChatMessageHistory
    for msg in msgs.messages:
        if msg.type == 'AIMessageChunk':
            st.chat_message('ai').write(msg.content)
        else:
            st.chat_message(msg.type).write(msg.content)

    # If user inputs a new prompt, generate and draw a new response
    if user_input := st.chat_input("How can I help?"):
        st.chat_message("human").write(user_input)

        # New messages are saved to history automatically by Langchain during run
        config = {"configurable": {"session_id": "any"}}
        st.chat_message('ai').write_stream(chain_with_history.stream({"question": user_input}, config))
    
    trigger_keywords = ["ppt", "powerpoint", "create ppt", "generate ppt", "presentation"]
    if user_input and any(keyword in user_input.lower() for keyword in trigger_keywords):
        try:
            pptx_path = generate_styled_pptx()
            with open(pptx_path, "rb") as f:
                st.download_button(
                    label="⬇️ Download Presentation (.pptx)",
                    data=f,
                    file_name="Styled_Business_Strategy_Proposal.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        except Exception as e:
            st.error(f"⚠️ Failed to generate presentation: {e}")