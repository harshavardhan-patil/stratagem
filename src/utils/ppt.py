from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
import tempfile
import random

ENABLE_DECORATIVE_CIRCLES = False
ENABLE_TITLE_BAR_PATTERNS = False
ENABLE_BACKGROUND_PATTERNS = True  # You can toggle this too

def get_rgb_components(color: RGBColor):
    color_str = str(color)  # e.g., 'FF6699'
    r = int(color_str[0:2], 16)
    g = int(color_str[2:4], 16)
    b = int(color_str[4:6], 16)
    return r, g, b

# --- Enhanced Style Constants ---
TITLE_FONT_SIZE = Pt(44)
SUBTITLE_FONT_SIZE = Pt(28)
HEADING_FONT_SIZE = Pt(32)
BODY_FONT_SIZE = Pt(20)

# Color schemes - you can switch between these
COLOR_SCHEMES = {
    "blue": {
        "primary": RGBColor(0, 102, 204),      # Main accent color
        "secondary": RGBColor(41, 128, 185),   # Secondary accent
        "accent": RGBColor(52, 152, 219),      # Highlight color
        "text": RGBColor(44, 62, 80),          # Text color
        "light_text": RGBColor(255, 255, 255), # Light text color
        "background": RGBColor(240, 240, 240), # Slide background
        "alt_background": RGBColor(245, 250, 255), # Alternative background
        "chart_colors": [
            RGBColor(41, 128, 185),
            RGBColor(26, 188, 156),
            RGBColor(46, 204, 113),
            RGBColor(241, 196, 15)
        ]
    },
    "green": {
        "primary": RGBColor(46, 139, 87),     
        "secondary": RGBColor(60, 179, 113),  
        "accent": RGBColor(76, 187, 23),      
        "text": RGBColor(44, 62, 80),         
        "light_text": RGBColor(255, 255, 255),
        "background": RGBColor(240, 245, 240),
        "alt_background": RGBColor(245, 255, 245),
        "chart_colors": [
            RGBColor(46, 139, 87),
            RGBColor(115, 169, 80),
            RGBColor(39, 174, 96),
            RGBColor(241, 196, 15)
        ]
    },
    "purple": {
        "primary": RGBColor(142, 68, 173),     
        "secondary": RGBColor(155, 89, 182),   
        "accent": RGBColor(187, 143, 206),      
        "text": RGBColor(44, 62, 80),          
        "light_text": RGBColor(255, 255, 255), 
        "background": RGBColor(245, 240, 250), 
        "alt_background": RGBColor(250, 245, 255),
        "chart_colors": [
            RGBColor(142, 68, 173),
            RGBColor(155, 89, 182),
            RGBColor(26, 188, 156),
            RGBColor(241, 196, 15)
        ]
    }
}

# Choose a color scheme (can be randomized or user-selected)
COLORS = COLOR_SCHEMES["blue"]  # Default color scheme

def clean_financial_slide(prs):
    """
    Removes non-text decorative shapes from the slide titled 'Financial Projections'.
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text is not None:
                if "financial projections" in shape.text.strip().lower():
                    shapes_to_remove = [
                        s for s in slide.shapes
                        if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and not s.has_text_frame
                    ]
                    for s in shapes_to_remove:
                        s._element.getparent().remove(s._element)
                    return
            

# --- Background Pattern Generators ---
def add_dotted_background(slide, color=None, density=120):
    """Add dotted pattern background to a slide"""
    if color is None:
        color = COLORS["accent"]
    
    # Create small dots scattered across the slide
    for _ in range(density):
        x = random.uniform(0, 10)
        y = random.uniform(0, 7.5)
        dot_size = random.uniform(0.02, 0.05)  # Small sized dots
        
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, 
            Inches(x), 
            Inches(y), 
            Inches(dot_size), 
            Inches(dot_size)
        )
        dot.fill.solid()
        
        # Create a slightly transparent version of the color
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        # Set transparency (0-100000, where 100000 is fully transparent)
        dot.fill.transparency = random.randint(60000, 85000)

def add_wave_background(slide, color1=None, color2=None):
    """Add a subtle wave pattern to slide background"""
    if color1 is None:
        color1 = COLORS["primary"]
    if color2 is None:
        color2 = COLORS["secondary"]
    
    # Create a few wave shapes
    for i in range(3):
        height = random.uniform(0.8, 1.2)
        y_pos = random.uniform(-0.5, 7)
        
        wave = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.WAVE, 
            Inches(-0.5), 
            Inches(y_pos), 
            Inches(11), 
            Inches(height)
        )
        wave.fill.solid()
        wave.fill.fore_color.rgb = color1 if i % 2 == 0 else color2
        wave.line.fill.background()
        wave.fill.transparency = random.randint(70000, 90000)

def add_geometric_background(slide):
    """Add geometric shapes as background elements"""
    shapes = [
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        MSO_AUTO_SHAPE_TYPE.OVAL,
        MSO_AUTO_SHAPE_TYPE.DIAMOND,
        MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM
        ]
    
    for _ in range(15):
        shape_type = random.choice(shapes)
        size = random.uniform(0.5, 2.0)
        x = random.uniform(-0.5, 10)
        y = random.uniform(-0.5, 7.5)
        
        shape = slide.shapes.add_shape(
            shape_type, 
            Inches(x), 
            Inches(y), 
            Inches(size), 
            Inches(size)
        )
        
        # Pick a color from the color scheme
        color_keys = ["primary", "secondary", "accent"]
        color = COLORS[random.choice(color_keys)]
        
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.fill.transparency = random.randint(80000, 95000)  # Very transparent

def add_gradient_background(slide, direction="vertical"):
    """
    This simulates a gradient effect by adding multiple semi-transparent rectangles
    """
    steps = 15  # Number of gradient steps
    width = Inches(10)
    height = Inches(7.5)
    
    # Choose two colors for the gradient
    start_color = COLORS["primary"]
    end_color = COLORS["secondary"]

    r1, g1, b1 = get_rgb_components(start_color)
    r2, g2, b2 = get_rgb_components(end_color)

    for i in range(steps):
        # Calculate interpolated color
        ratio = i / steps
        r = int(r1 * (1 - ratio) + r2 * ratio)
        g = int(g1 * (1 - ratio) + g2 * ratio)
        b = int(b1 * (1 - ratio) + b2 * ratio)
        
        if direction == "vertical":
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0),
                Inches(7.5 * i / steps),
                width,
                Inches(7.5 / steps + 0.1)  # Slight overlap to avoid gaps
            )
        else:  # horizontal
            rect = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(10 * i / steps),
                Inches(0),
                Inches(10 / steps + 0.1),  # Slight overlap
                height
            )
            
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.fill.background()
        rect.fill.transparency = 50000  # Semi-transparent

# --- Title Slide ---
def add_styled_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    add_geometric_background(slide)

    center_overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    )
    center_overlay.fill.solid()
    center_overlay.fill.fore_color.rgb = COLORS["background"]
    center_overlay.line.fill.background()
    center_overlay.fill.transparency = 20000

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.8), Inches(2), Inches(0.15), Inches(3.5)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS["primary"]
    accent_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(2.1), Inches(7.4), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS["text"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    subtitle_box = slide.shapes.add_textbox(Inches(1.6), Inches(3.5), Inches(7.4), Inches(1))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    p = subtitle_tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS["secondary"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    if ENABLE_DECORATIVE_CIRCLES:
        for i in range(3):
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(9 - i*0.3), Inches(5.5 + i*0.2), Inches(0.4), Inches(0.4)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS["primary"] if i % 2 == 0 else COLORS["accent"]
            circle.line.fill.background()

# --- Section Divider Slide ---
def add_section_divider_slide(prs, title):
    """Creates an impactful section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add gradient background
    add_gradient_background(slide, direction="horizontal")
    
    # Add a semi-transparent overlay for better text visibility
    overlay = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(10),
        Inches(7.5)
    )
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS["primary"]
    overlay.line.fill.background()
    overlay.fill.transparency = 40000  # Somewhat transparent
    
    # Large, bold title text
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = title.upper()  # Uppercase for impact
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    
    # Add decorative line
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(4),
        Inches(4.5),
        Inches(2),
        Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["light_text"]
    line.line.fill.background()

# --- Enhanced Bullet Slide with Optional Infographic ---
def add_bullet_slide(prs, title, bullets, infographic=None, background_style="dots"):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set base background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    
    # Add background pattern based on style
    if ENABLE_BACKGROUND_PATTERNS:
        if background_style == "dots":
            add_dotted_background(slide)
        elif background_style == "waves":
            add_wave_background(slide)
        elif background_style == "geometric":
            add_geometric_background(slide)
    
    # Enhanced title bar
    title_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 
        Inches(0), 
        Inches(0), 
        Inches(10), 
        Inches(1.2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS["primary"]
    title_shape.line.fill.background()
    
    # Add subtle pattern to title bar
    if ENABLE_TITLE_BAR_PATTERNS:
        for i in range(5):
            x = random.uniform(0, 9)
            circle = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(x),
                Inches(random.uniform(0.1, 0.9)),
                Inches(0.3),
                Inches(0.3)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = COLORS["secondary"]
            circle.line.fill.background()
            circle.fill.transparency = 60000
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    p.text = title
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = COLORS["light_text"]
    p.font.name = "Calibri"
    
    # Content area with slight shadow effect
    content_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(1.3),
        Inches(6.8),
        Inches(5.5)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_bg.line.color.rgb = RGBColor(220, 220, 220)
    
    # Content with custom bullets
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(6.5), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        para = tf.add_paragraph() if i != 0 else tf.paragraphs[0]
        
        # Alternate bullet styles for visual interest
        if i % 2 == 0:
            para.text = f"► {bullet}"  # Triangle bullet
        else:
            para.text = f"• {bullet}"  # Round bullet
            
        para.font.size = BODY_FONT_SIZE
        para.font.color.rgb = COLORS["text"]
        para.font.name = "Calibri"
        
        # Add spacing between bullets
        if i > 0:
            para.space_before = Pt(12)

    # Add infographic if provided
    if infographic:
        infographic(slide)

# --- Enhanced SWOT Infographic Slide ---
def add_swot_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    add_dotted_background(slide)

    tb = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = COLORS["primary"]
    tb.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "SWOT Analysis"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["light_text"]
    title_p.font.name = "Calibri"

    labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    coords = [(0.6, 1.4), (5.1, 1.4), (0.6, 4.0), (5.1, 4.0)]
    colors = [
        RGBColor(76, 175, 80),
        RGBColor(244, 67, 54),
        RGBColor(33, 150, 243),
        RGBColor(255, 152, 0)
    ]

    for i, label in enumerate(labels):
        x, y = coords[i]

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(2.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = colors[i]
        box.line.width = Pt(2.5)

        header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors[i]
        header.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.2), Inches(0.5))
        label_tf = label_box.text_frame
        p = label_tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        text_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.7), Inches(3.8), Inches(1.4)
        )
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        placeholder_p = text_tf.paragraphs[0]
        placeholder_p.text = f"Add {label.lower()} here"
        placeholder_p.font.italic = True
        placeholder_p.font.size = Pt(14)
        placeholder_p.font.color.rgb = RGBColor(120, 120, 120)

# --- Roadmap Slide ---
def add_roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    add_dotted_background(slide)

    tb = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = COLORS["primary"]
    tb.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "Implementation Roadmap"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["light_text"]
    title_p.font.name = "Calibri"

    timeline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(3), Inches(8), Inches(0.1)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = COLORS["text"]
    timeline.line.fill.background()

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(9), Inches(2.9), Inches(0.5), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["text"]
    arrow.line.fill.background()

    phases = [
        {"name": "Planning", "duration": "0-3 months", "color": COLORS["chart_colors"][0]},
        {"name": "Rollout", "duration": "4-6 months", "color": COLORS["chart_colors"][1]},
        {"name": "Expansion", "duration": "7-12 months", "color": COLORS["chart_colors"][2]}
    ]

    for i, phase in enumerate(phases):
        x = 2 + i * 3

        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x), Inches(3), Inches(0.8), Inches(0.8)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = phase["color"]
        node.line.fill.background()
        p = node.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["light_text"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        lbl = slide.shapes.add_textbox(
            Inches(x - 0.6), Inches(4), Inches(2), Inches(0.6)
        )
        p2 = lbl.text_frame.paragraphs[0]
        p2.text = phase["name"]
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLORS["text"]
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER

        dur = slide.shapes.add_textbox(
            Inches(x - 0.6), Inches(4.4), Inches(2), Inches(0.5)
        )
        p3 = dur.text_frame.paragraphs[0]
        p3.text = phase["duration"]
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = COLORS["secondary"]
        p3.font.name = "Calibri"
        p3.alignment = PP_ALIGN.CENTER

        act_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x - 0.75), Inches(1.8), Inches(2.3), Inches(0.9)
        )
        act_box.fill.solid()
        act_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        act_box.line.color.rgb = phase["color"]

        act_text = slide.shapes.add_textbox(
            Inches(x - 0.65), Inches(1.9), Inches(2.1), Inches(0.7)
        )
        act_tf = act_text.text_frame
        act_tf.word_wrap = True
        p4 = act_tf.paragraphs[0]
        p4.text = "Key Activities:"
        p4.font.size = Pt(12)
        p4.font.bold = True
        p4.font.color.rgb = COLORS["text"]
        p5 = act_tf.add_paragraph()
        p5.text = "• Activity 1\n• Activity 2"
        p5.font.size = Pt(10)
        p5.font.color.rgb = COLORS["text"]

# --- Financial Chart Slide ---
def add_financial_chart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "Financial Projections"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["primary"]
    title_p.font.name = "Calibri"

    # Chart bars (simple visual simulation)
    items = [
        ("Initial Investment", 2.5, COLORS["chart_colors"][0]),
        ("Expected Revenue", 4.2, COLORS["chart_colors"][1]),
        ("Operational Costs", 2.0, COLORS["chart_colors"][2]),
        ("Break-even", 3.1, COLORS["chart_colors"][3])
    ]

    for i, (label, height, color) in enumerate(items):
        x = 1.0 + i * 2.1
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(5.5 - height),
            Inches(1),
            Inches(height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(5.7), Inches(1.4), Inches(0.5))
        label_tf = label_box.text_frame
        p = label_tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["text"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

def add_swot_slide_with_data(prs, swot_data):
    """Creates a SWOT analysis slide with the provided data"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    if ENABLE_BACKGROUND_PATTERNS:
        add_dotted_background(slide)

    tb = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = COLORS["primary"]
    tb.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "SWOT Analysis"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["light_text"]
    title_p.font.name = "Calibri"

    labels = ["Strengths", "Weaknesses", "Opportunities", "Threats"]
    categories = ["strengths", "weaknesses", "opportunities", "threats"]
    coords = [(0.6, 1.4), (5.1, 1.4), (0.6, 4.0), (5.1, 4.0)]
    colors = [
        RGBColor(76, 175, 80),
        RGBColor(244, 67, 54),
        RGBColor(33, 150, 243),
        RGBColor(255, 152, 0)
    ]

    for i, (label, category) in enumerate(zip(labels, categories)):
        x, y = coords[i]

        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(2.3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = colors[i]
        box.line.width = Pt(2.5)

        header = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x), Inches(y), Inches(4.2), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = colors[i]
        header.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.2), Inches(0.5))
        label_tf = label_box.text_frame
        p = label_tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        text_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y + 0.7), Inches(3.8), Inches(1.4)
        )
        text_tf = text_box.text_frame
        text_tf.word_wrap = True
        
        # Add actual content from swot_data
        items = swot_data[category][:3]  # Limit to top 3 for space reasons
        
        for j, item in enumerate(items):
            if j == 0:
                p = text_tf.paragraphs[0]
            else:
                p = text_tf.add_paragraph()
                
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS["text"]
            p.space_after = Pt(6)

    return slide

def add_roadmap_slide_with_data(prs, roadmap_data):
    """Creates a roadmap slide with the provided timeline data"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]
    if ENABLE_BACKGROUND_PATTERNS:
        add_dotted_background(slide)

    tb = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(1.2)
    )
    tb.fill.solid()
    tb.fill.fore_color.rgb = COLORS["primary"]
    tb.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "Implementation Roadmap"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["light_text"]
    title_p.font.name = "Calibri"

    timeline = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1), Inches(3), Inches(8), Inches(0.1)
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = COLORS["text"]
    timeline.line.fill.background()

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(9), Inches(2.9), Inches(0.5), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLORS["text"]
    arrow.line.fill.background()

    phases = [
        {"name": "Planning", "duration": "0-3 months", "activities": roadmap_data["short_term"], "color": COLORS["chart_colors"][0]},
        {"name": "Rollout", "duration": "4-6 months", "activities": roadmap_data["mid_term"], "color": COLORS["chart_colors"][1]},
        {"name": "Expansion", "duration": "7-12 months", "activities": roadmap_data["long_term"], "color": COLORS["chart_colors"][2]}
    ]

    for i, phase in enumerate(phases):
        x = 2 + i * 3

        node = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x), Inches(3), Inches(0.8), Inches(0.8)
        )
        node.fill.solid()
        node.fill.fore_color.rgb = phase["color"]
        node.line.fill.background()
        p = node.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS["light_text"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        lbl = slide.shapes.add_textbox(
            Inches(x - 0.6), Inches(4), Inches(2), Inches(0.6)
        )
        p2 = lbl.text_frame.paragraphs[0]
        p2.text = phase["name"]
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLORS["text"]
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER

        dur = slide.shapes.add_textbox(
            Inches(x - 0.6), Inches(4.4), Inches(2), Inches(0.5)
        )
        p3 = dur.text_frame.paragraphs[0]
        p3.text = phase["duration"]
        p3.font.size = Pt(14)
        p3.font.italic = True
        p3.font.color.rgb = COLORS["secondary"]
        p3.font.name = "Calibri"
        p3.alignment = PP_ALIGN.CENTER

        act_box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x - 0.75), Inches(1.8), Inches(2.3), Inches(0.9)
        )
        act_box.fill.solid()
        act_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        act_box.line.color.rgb = phase["color"]

        act_text = slide.shapes.add_textbox(
            Inches(x - 0.65), Inches(1.9), Inches(2.1), Inches(0.7)
        )
        act_tf = act_text.text_frame
        act_tf.word_wrap = True
        p4 = act_tf.paragraphs[0]
        p4.text = "Key Activities:"
        p4.font.size = Pt(12)
        p4.font.bold = True
        p4.font.color.rgb = COLORS["text"]
        
        # Add actual activities from the data
        activities_text = ""
        for activity in phase["activities"][:2]:  # Limit to top 2 activities for space
            activities_text += f"• {activity}\n"
            
        p5 = act_tf.add_paragraph()
        p5.text = activities_text.strip()
        p5.font.size = Pt(10)
        p5.font.color.rgb = COLORS["text"]

    return slide

def add_financial_chart_slide_with_data(prs, financial_data):
    """Creates a financial chart slide with the provided data"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["background"]

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "Financial Projections"
    title_p.font.size = TITLE_FONT_SIZE
    title_p.font.bold = True
    title_p.font.color.rgb = COLORS["primary"]
    title_p.font.name = "Calibri"

    # Extract key metrics from financial data
    metrics = []
    for item in financial_data:
        if "revenue" in item.lower():
            metrics.append(("Revenue", 4.2, COLORS["chart_colors"][0]))
        elif "investment" in item.lower():
            metrics.append(("Investment", 2.5, COLORS["chart_colors"][1]))
        elif "cost" in item.lower() or "expense" in item.lower():
            metrics.append(("Costs", 2.0, COLORS["chart_colors"][2]))
        elif "break" in item.lower() or "profit" in item.lower():
            metrics.append(("Profit", 3.1, COLORS["chart_colors"][3]))
    
    # If we didn't extract enough metrics, add defaults
    default_items = [
        ("Revenue", 4.2, COLORS["chart_colors"][0]),
        ("Investment", 2.5, COLORS["chart_colors"][1]),
        ("Costs", 2.0, COLORS["chart_colors"][2]),
        ("Profit", 3.1, COLORS["chart_colors"][3])
    ]
    
    # Use extracted metrics, fill with defaults if needed
    items = metrics[:4] if metrics else default_items
    
    # Add the financial data as text
    data_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.3))
    data_tf = data_box.text_frame
    data_tf.word_wrap = True
    
    for i, item in enumerate(financial_data[:3]):  # Top 3 financial points
        if i == 0:
            p = data_tf.paragraphs[0]
        else:
            p = data_tf.add_paragraph()
            
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS["text"]
        p.space_after = Pt(6)

    # Chart bars (simple visual simulation)
    for i, (label, height, color) in enumerate(items):
        x = 1.0 + i * 2.1
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(5.5 - height),
            Inches(1),
            Inches(height)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x - 0.2), Inches(5.7), Inches(1.4), Inches(0.5))
        label_tf = label_box.text_frame
        p = label_tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["text"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER



# --- Generate Full Styled PPT ---
def generate_styled_pptx():
    prs = Presentation()

    add_styled_title_slide(prs, "Business Growth Strategy Proposal", "Prepared by Strategic Synthesis AI")

    add_section_divider_slide(prs, "Introduction")
    add_bullet_slide(prs, "Company Overview", [
        "Founded in 2020 with a vision to democratize access to AI.",
        "Headquartered in New York with global outreach.",
        "Team of 50+ AI experts and software engineers."
    ])

    add_section_divider_slide(prs, "Business Landscape")
    add_bullet_slide(prs, "Market Understanding", [
        "Target market includes mid-sized B2B SaaS providers.",
        "Market size is projected to reach $5B by 2027.",
        "Customer pain points revolve around data silos and integration."
    ])
    add_swot_slide(prs)

    add_section_divider_slide(prs, "Strategy & Roadmap")
    add_bullet_slide(prs, "Strategic Objectives", [
        "Expand to 3 new international markets.",
        "Achieve $10M ARR by Q4 2026.",
        "Establish partnerships with key ecosystem players."
    ])
    add_roadmap_slide(prs)

    add_section_divider_slide(prs, "Financial Outlook")
    add_financial_chart_slide(prs)
    clean_financial_slide(prs)

    add_section_divider_slide(prs, "Outcomes & Next Steps")
    add_bullet_slide(prs, "Expected Impact", [
        "Revenue growth by 300% over 2 years.",
        "Improved brand positioning in AI infrastructure sector.",
        "Sustainable operational margin with lean team."
    ])

    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(output_file.name)
    return output_file.name